"""Document-to-text extraction for ``read_file``.

Jupyter notebooks, DOCX, XLSX, PPTX, RTF, and OpenDocument (ODT/ODS/ODP) are
parsed with the stdlib alone (no hard dependencies). Legacy binary ``.doc``
is also stdlib-only, via a best-effort/lossy byte-scrape (no OLE2 stream
parsing — see ``_extract_legacy_doc``).

PDF and image extraction additionally need ``PyMuPDF`` (``fitz``) and
``pytesseract`` plus a ``tesseract`` binary on PATH (installed by
scripts/install.ps1 / install.sh); PDF also has a lightweight ``pypdf``
fallback for the rare file PyMuPDF fails to open (text-layer only, no OCR).
All of these import lazily and raise :class:`ExtractionError` when
unavailable, same as any other malformed document.

Malformed documents raise :class:`ExtractionError`; callers can then fall back to
normal text/binary handling.
"""

from __future__ import annotations

import json
import logging
import posixpath
import re
import shutil
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

__all__ = ["EXTRACTABLE_EXTENSIONS", "ExtractionError", "extract_document_text", "is_extractable_document"]

logger = logging.getLogger(__name__)

_IMAGE_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"})
_ODF_EXTENSIONS = frozenset({".odt", ".ods", ".odp"})
_HTML_EXTENSIONS = frozenset({".html", ".htm"})
EXTRACTABLE_EXTENSIONS = (
    frozenset({".ipynb", ".docx", ".xlsx", ".pptx", ".pdf", ".rtf", ".doc", ".msg", ".eml", ".epub"})
    | _IMAGE_EXTENSIONS
    | _ODF_EXTENSIONS
    | _HTML_EXTENSIONS
)
MAX_XLSX_BYTES = 50 * 1024 * 1024
_MAX_XLSX_ROWS_PER_SHEET = 5000
_MAX_XLSX_COLS = 256

_NS_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_NS_S = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_NS_ODF_TEXT = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"


class ExtractionError(Exception):
    """Raised when a supported-looking document cannot be rendered as text."""


def _extension(path: str) -> str:
    ext = Path(path).suffix.lower()
    return ext if ext in EXTRACTABLE_EXTENSIONS else ""


def is_extractable_document(path: str) -> bool:
    return bool(_extension(path))


def _first_non_empty(path: str, extractors, label: str) -> str:
    """Run *extractors* in order; return the first non-empty result.

    Each extractor is a ``(name, callable)`` pair taking *path*. Any failure --
    a missing optional library, a parser that chokes on this particular file --
    moves on to the next candidate, so a document only fails once every
    strategy has. The final :class:`ExtractionError` names what was tried, since
    "PDF contains no extractable text" is useless when the real cause is that
    three of four extractors were never installed.
    """
    attempts: list[str] = []
    for name, extractor in extractors:
        try:
            text = extractor(path)
        except ExtractionError as exc:
            attempts.append(f"{name}: {exc}")
            continue
        except Exception as exc:  # noqa: BLE001 - a broken parser must not mask the rest
            logger.debug("%s extractor failed for %s", name, path, exc_info=True)
            attempts.append(f"{name}: {exc}")
            continue
        if text and text.strip():
            return text
        attempts.append(f"{name}: no text")
    raise ExtractionError(f"Could not extract text from this {label}. Tried -- " + "; ".join(attempts))


def extract_document_text(path: str) -> str:
    ext = _extension(path)
    if ext == ".ipynb":
        return _extract_notebook(path)
    if ext == ".docx":
        return _first_non_empty(
            path,
            [("stdlib", _extract_docx), ("python-docx", _extract_docx_lib)],
            "Word document",
        )
    if ext == ".xlsx":
        return _first_non_empty(
            path,
            [("stdlib", _extract_xlsx), ("openpyxl", _extract_xlsx_lib)],
            "spreadsheet",
        )
    if ext == ".pptx":
        return _first_non_empty(
            path,
            [("stdlib", _extract_pptx), ("python-pptx", _extract_pptx_lib)],
            "presentation",
        )
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".rtf":
        return _first_non_empty(
            path,
            [("stdlib", _extract_rtf), ("striprtf", _extract_rtf_lib)],
            "RTF document",
        )
    if ext == ".doc":
        return _first_non_empty(
            path,
            [("olefile", _extract_legacy_doc_olefile), ("byte-scrape", _extract_legacy_doc)],
            "legacy Word document",
        )
    if ext == ".msg":
        return _extract_msg(path)
    if ext == ".eml":
        return _extract_eml(path)
    if ext == ".epub":
        return _extract_epub(path)
    if ext in _HTML_EXTENSIONS:
        return _extract_html(path)
    if ext in _ODF_EXTENSIONS:
        return _first_non_empty(
            path,
            [("stdlib", _extract_odf), ("odfpy", _extract_odf_lib)],
            "OpenDocument file",
        )
    if ext in _IMAGE_EXTENSIONS:
        return _extract_image(path)
    raise ExtractionError(f"Unsupported document type: {path!r}")


def _source_text(source) -> str:
    if isinstance(source, str):
        return source
    if isinstance(source, list):
        return "".join(item for item in source if isinstance(item, str))
    return ""


def _extract_notebook(path: str) -> str:
    try:
        with open(path, encoding="utf-8", errors="replace") as fh:
            nb = json.load(fh)
    except (OSError, ValueError, json.JSONDecodeError) as exc:
        raise ExtractionError(f"Not a valid notebook: {exc}") from exc
    if not isinstance(nb, dict):
        raise ExtractionError("Notebook root is not an object")

    cells = nb.get("cells")
    if not isinstance(cells, list):
        cells = [
            cell
            for ws in nb.get("worksheets", [])
            if isinstance(ws, dict)
            for cell in ws.get("cells", [])
        ]
    if not cells:
        raise ExtractionError("Notebook contains no cells")

    counts = {"markdown": 0, "code": 0, "raw": 0}
    labels = {"markdown": "Markdown", "code": "Code", "raw": "Raw"}
    out: list[str] = []
    for cell in cells:
        if not isinstance(cell, dict):
            continue
        typ = cell.get("cell_type")
        if typ not in labels:
            continue
        counts[typ] += 1
        suffix = f" {counts[typ]}" if typ != "raw" else ""
        out.extend((f"# ── {labels[typ]} cell{suffix} ──", _source_text(cell.get("source", "")).rstrip("\n"), ""))
    if not out:
        raise ExtractionError("Notebook contains no readable cells")
    return "\n".join(out).rstrip("\n") + "\n"


def _zip_xml(zf: zipfile.ZipFile, name: str) -> ET.Element:
    try:
        return ET.fromstring(zf.read(name))
    except KeyError as exc:
        raise ExtractionError(f"Missing {name}") from exc
    except ET.ParseError as exc:
        raise ExtractionError(f"Malformed XML in {name}: {exc}") from exc


def _extract_docx(path: str) -> str:
    try:
        with zipfile.ZipFile(path) as zf:
            root = _zip_xml(zf, "word/document.xml")
    except zipfile.BadZipFile as exc:
        raise ExtractionError(f"Not a valid DOCX: {exc}") from exc
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc

    w = f"{{{_NS_W}}}"
    lines: list[str] = []
    for para in root.iter(f"{w}p"):
        buf: list[str] = []
        for node in para.iter():
            if node.tag == f"{w}t":
                buf.append(node.text or "")
            elif node.tag == f"{w}tab":
                buf.append("\t")
            elif node.tag in {f"{w}br", f"{w}cr"}:
                buf.append("\n")
        lines.extend("".join(buf).split("\n"))
    if not any(line.strip() for line in lines):
        raise ExtractionError("DOCX contains no extractable text")
    return "\n".join(lines).rstrip("\n") + "\n"


def _extract_xlsx(path: str) -> str:
    try:
        with zipfile.ZipFile(path) as zf:
            names = set(zf.namelist())
            shared = _shared_strings(zf, names)
            sheets = _workbook_sheets(zf)
            rels = _workbook_rels(zf, names)
            out: list[str] = []
            for name, state, rid in sheets:
                if state in {"hidden", "veryHidden"}:
                    continue
                part = _sheet_part(rels.get(rid, ""))
                if part not in names:
                    continue
                try:
                    rows = _sheet_rows(zf.read(part), shared)
                except ET.ParseError:
                    continue
                out.append(f"# ── Sheet: {name} ──")
                out.extend("\t".join(row) for row in rows)
                if not rows:
                    out.append("(empty)")
                out.append("")
    except zipfile.BadZipFile as exc:
        raise ExtractionError(f"Not a valid XLSX: {exc}") from exc
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc

    if not out:
        raise ExtractionError("XLSX has no visible sheets with content")
    return "\n".join(out).rstrip("\n") + "\n"


def _extract_pptx(path: str) -> str:
    try:
        with zipfile.ZipFile(path) as zf:
            names = zf.namelist()
            slide_names = sorted(
                (n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)),
                key=lambda n: int(re.search(r"\d+", n).group()),
            )
            a = f"{{{_NS_A}}}"
            out: list[str] = []
            for idx, name in enumerate(slide_names, start=1):
                try:
                    root = ET.fromstring(zf.read(name))
                except ET.ParseError:
                    continue
                out.append(f"# ── Slide {idx} ──")
                for para in root.iter(f"{a}p"):
                    out.append("".join(t.text or "" for t in para.iter(f"{a}t")))
                out.append("")
    except zipfile.BadZipFile as exc:
        raise ExtractionError(f"Not a valid PPTX: {exc}") from exc
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc

    if not any(line.strip() and not line.startswith("#") for line in out):
        raise ExtractionError("PPTX contains no extractable text")
    return "\n".join(out).rstrip("\n") + "\n"


def _extract_odf(path: str) -> str:
    """OpenDocument Text/Spreadsheet/Presentation (.odt/.ods/.odp).

    All three share the same content.xml structure closely enough that a
    single text-node walk (paragraphs + headings) covers them.
    """
    try:
        with zipfile.ZipFile(path) as zf:
            root = _zip_xml(zf, "content.xml")
    except zipfile.BadZipFile as exc:
        raise ExtractionError(f"Not a valid OpenDocument file: {exc}") from exc
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc

    t = f"{{{_NS_ODF_TEXT}}}"
    lines = ["".join(node.itertext()) for node in root.iter() if node.tag in (f"{t}p", f"{t}h")]
    if not any(line.strip() for line in lines):
        raise ExtractionError("Document contains no extractable text")
    return "\n".join(lines).rstrip("\n") + "\n"


# Destination groups (``{\fonttbl ...}`` etc.) hold document metadata, not
# body text -- their control word never yields visible text but the plain
# characters inside the group (e.g. a font name) do, so they must be tracked
# and skipped explicitly rather than just stripping control words.
_RTF_SKIP_DESTINATIONS = frozenset({
    "fonttbl", "colortbl", "stylesheet", "info", "generator", "pict",
    "object", "footer", "footerf", "footerl", "footerr", "header",
    "headerf", "headerl", "headerr", "footnote", "themedata",
    "colorschememapping", "latentstyles", "rsid", "xmlnstbl", "listtable",
    "listoverridetable", "revtbl", "datastore", "companyname", "operator",
    "template", "doccomm", "keywords", "subject", "title", "author",
})
_RTF_HEX_ESCAPE = re.compile(r"\\'([0-9a-fA-F]{2})")
_RTF_CONTROL_WORD = re.compile(r"\\([a-zA-Z]+)(-?\d+)? ?")
_RTF_ESCAPED_LITERAL = re.compile(r"\\([{}\\])")


def _extract_rtf(path: str) -> str:
    """Stdlib-only RTF extraction: a small brace-tracking parser rather than
    a flat regex strip, so destination groups (font/color tables, generator,
    embedded objects, ``\\*`` optional destinations, ...) are skipped
    wholesale instead of leaking their plain-text contents (a font name,
    etc.) into the output. Lossy on ``\\uNNNN`` unicode escapes and complex
    field codes, but recovers clean body text without a dependency like
    ``striprtf``.
    """
    try:
        raw = Path(path).read_bytes()
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc
    if not raw.lstrip().startswith(b"{\\rtf"):
        raise ExtractionError("Not a valid RTF file")
    text = raw.decode("latin-1", errors="replace")

    out: list[str] = []
    skip_stack = [False]
    i, n = 0, len(text)
    while i < n:
        ch = text[i]
        if ch == "{":
            skip_stack.append(skip_stack[-1])
            i += 1
        elif ch == "}":
            if len(skip_stack) > 1:
                skip_stack.pop()
            i += 1
        elif ch == "\\":
            if m := _RTF_HEX_ESCAPE.match(text, i):
                if not skip_stack[-1]:
                    out.append(bytes([int(m.group(1), 16)]).decode("latin-1", errors="replace"))
                i = m.end()
            elif m := _RTF_CONTROL_WORD.match(text, i):
                word = m.group(1)
                if word in ("par", "line"):
                    if not skip_stack[-1]:
                        out.append("\n")
                elif word == "tab":
                    if not skip_stack[-1]:
                        out.append("\t")
                elif word in _RTF_SKIP_DESTINATIONS:
                    skip_stack[-1] = True
                i = m.end()
            elif text.startswith("\\*", i):
                skip_stack[-1] = True  # unknown/optional destination — skip
                i += 2
            elif m := _RTF_ESCAPED_LITERAL.match(text, i):
                if not skip_stack[-1]:
                    out.append(m.group(1))
                i = m.end()
            else:
                i += 2  # unrecognized control symbol, e.g. \~
        else:
            if not skip_stack[-1]:
                out.append(ch)
            i += 1

    stripped = re.sub(r"[ \t]*\n[ \t]*", "\n", "".join(out)).strip(" \t\n")
    stripped = re.sub(r"\n{3,}", "\n\n", stripped)
    if not stripped:
        raise ExtractionError("RTF contains no extractable text")
    return stripped + "\n"


_DOC_OLE2_SIGNATURE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_DOC_MIN_RUN = 4


def _printable_utf16_runs(data: bytes) -> str:
    """Pull runs of printable UTF-16LE characters out of raw bytes.

    Shared by both legacy-.doc strategies: ``_extract_legacy_doc_olefile``
    feeds it just the WordDocument stream, ``_extract_legacy_doc`` feeds it
    the whole file. Neither walks the piece table, so both lose paragraph
    structure -- this only recovers the characters.
    """
    runs: list[str] = []
    current: list[str] = []
    for i in range(0, len(data) - 1, 2):
        lo, hi = data[i], data[i + 1]
        if hi == 0 and 0x20 <= lo < 0x7F:
            current.append(chr(lo))
            continue
        if len(current) >= _DOC_MIN_RUN:
            runs.append("".join(current))
        current = []
    if len(current) >= _DOC_MIN_RUN:
        runs.append("".join(current))

    if not runs:
        raise ExtractionError("No extractable text found in legacy .doc file")
    return "\n".join(runs).rstrip("\n") + "\n"


def _extract_legacy_doc(path: str) -> str:
    """Last-resort text scrape for legacy binary (OLE2) ``.doc`` files.

    Scans the ENTIRE file for printable UTF-16LE runs — no OLE2 stream
    parsing at all, so it can pick up stray text from unrelated streams.
    Only reached when the olefile strategy is unavailable or fails; it needs
    no dependency, which is the whole point of keeping it.
    """
    try:
        data = Path(path).read_bytes()
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc
    if data[:8] != _DOC_OLE2_SIGNATURE:
        raise ExtractionError("Not a legacy OLE2 .doc file")
    return _printable_utf16_runs(data)


# ── Library-backed fallbacks ────────────────────────────────────────────
# Each mirrors a stdlib parser above and is only reached via _first_non_empty
# when that parser came back empty (tables, text boxes, charts, odd encodings)
# or raised. All imports are local so a missing lib is just one failed
# candidate in the chain, never an import error at module load.


def _require(module: str, pip_name: str):
    try:
        return __import__(module)
    except ImportError as exc:
        raise ExtractionError(f"{pip_name} is not installed: {exc}") from exc


def _extract_docx_lib(path: str) -> str:
    _require("docx", "python-docx")
    from docx import Document

    doc = Document(path)
    lines = [para.text for para in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            lines.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(lines).rstrip("\n") + "\n"


def _extract_pptx_lib(path: str) -> str:
    _require("pptx", "python-pptx")
    from pptx import Presentation

    out: list[str] = []
    for idx, slide in enumerate(Presentation(path).slides, start=1):
        out.append(f"# ── Slide {idx} ──")
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    out.append("\t".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append(f"(speaker notes) {notes}")
        out.append("")
    return "\n".join(out).rstrip("\n") + "\n"


def _extract_xlsx_lib(path: str) -> str:
    _require("openpyxl", "openpyxl")
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        out: list[str] = []
        for sheet in wb.worksheets:
            if sheet.sheet_state != "visible":
                continue
            # Cap at the sheet's real width, not _MAX_XLSX_COLS: passing the cap
            # as max_col makes openpyxl pad every row out to 256 columns, which
            # is a wall of tabs in the model's context for a 2-column sheet.
            max_col = min(sheet.max_column or 1, _MAX_XLSX_COLS)
            rows = [
                ["" if v is None else str(v) for v in row]
                for row in sheet.iter_rows(
                    max_row=_MAX_XLSX_ROWS_PER_SHEET, max_col=max_col, values_only=True
                )
            ]
            # Trim trailing blank cells/rows, matching the stdlib parser.
            for row in rows:
                while row and not row[-1].strip():
                    row.pop()
            while rows and not any(cell.strip() for cell in rows[-1]):
                rows.pop()
            out.append(f"# ── Sheet: {sheet.title} ──")
            out.extend("\t".join(row) for row in rows)
            if not rows:
                out.append("(empty)")
            out.append("")
    finally:
        wb.close()
    return "\n".join(out).rstrip("\n") + "\n"


def _extract_odf_lib(path: str) -> str:
    _require("odf", "odfpy")
    from odf import teletype, text
    from odf.opendocument import load

    doc = load(path)
    nodes = list(doc.getElementsByType(text.P)) + list(doc.getElementsByType(text.H))
    return "\n".join(teletype.extractText(node) for node in nodes).rstrip("\n") + "\n"


def _extract_rtf_lib(path: str) -> str:
    _require("striprtf", "striprtf")
    from striprtf.striprtf import rtf_to_text

    raw = Path(path).read_text(encoding="latin-1", errors="replace")
    # striprtf is permissive: hand it a non-RTF file and it hands the raw bytes
    # straight back as "text". Gate on the signature like the stdlib parser
    # does, so a mis-named file fails instead of yielding plausible garbage.
    if not raw.lstrip().startswith("{\\rtf"):
        raise ExtractionError("Not a valid RTF file")
    return rtf_to_text(raw, errors="ignore").rstrip("\n") + "\n"


def _extract_legacy_doc_olefile(path: str) -> str:
    """Legacy binary .doc via real OLE2 stream access.

    Still not a full piece-table/FIB walk (that's what antiword/LibreOffice
    are for) -- but reading the WordDocument stream directly beats scanning
    the whole file, because it can't pick up text from unrelated streams the
    way the byte-scrape fallback does.
    """
    olefile = _require("olefile", "olefile")
    if not olefile.isOleFile(path):
        raise ExtractionError("Not a legacy OLE2 .doc file")
    ole = olefile.OleFileIO(path)
    try:
        if not ole.exists("WordDocument"):
            raise ExtractionError("OLE2 file has no WordDocument stream (not a Word .doc?)")
        data = ole.openstream("WordDocument").read()
    finally:
        ole.close()
    return _printable_utf16_runs(data)


# ── Formats with no stdlib parser ───────────────────────────────────────


def _extract_msg(path: str) -> str:
    _require("extract_msg", "extract-msg")
    import extract_msg

    # extract_msg raises its own exception tree (InvalidFileFormatError, ...).
    # Translate to ExtractionError -- this module's contract is that a
    # malformed document raises ExtractionError so read_file can fall back to
    # binary/text handling. A leaked library exception breaks that.
    try:
        msg = extract_msg.Message(path)
    except Exception as exc:  # noqa: BLE001 - third-party exception tree
        raise ExtractionError(f"Not a readable Outlook .msg file: {exc}") from exc
    try:
        header = [
            f"From: {msg.sender or ''}",
            f"To: {msg.to or ''}",
            f"Cc: {msg.cc or ''}",
            f"Date: {msg.date or ''}",
            f"Subject: {msg.subject or ''}",
        ]
        attachments = [a.longFilename or a.shortFilename or "(unnamed)" for a in (msg.attachments or [])]
        if attachments:
            header.append(f"Attachments: {', '.join(str(a) for a in attachments)}")
        body = msg.body or ""
    except Exception as exc:  # noqa: BLE001 - third-party exception tree
        raise ExtractionError(f"Could not read Outlook .msg contents: {exc}") from exc
    finally:
        msg.close()
    if not body.strip():
        raise ExtractionError("Outlook message has no text body")
    return "\n".join(header) + "\n\n" + body.rstrip("\n") + "\n"


def _extract_eml(path: str) -> str:
    """RFC-822 email -- stdlib ``email``, no dependency."""
    from email import policy
    from email.parser import BytesParser

    try:
        with open(path, "rb") as fh:
            msg = BytesParser(policy=policy.default).parse(fh)
    except (OSError, ValueError) as exc:
        raise ExtractionError(f"Not a readable .eml file: {exc}") from exc

    header = [
        f"From: {msg.get('From', '')}",
        f"To: {msg.get('To', '')}",
        f"Cc: {msg.get('Cc', '')}",
        f"Date: {msg.get('Date', '')}",
        f"Subject: {msg.get('Subject', '')}",
    ]
    body_part = msg.get_body(preferencelist=("plain", "html"))
    if body_part is None:
        raise ExtractionError("Email has no text body")
    body = body_part.get_content()
    if body_part.get_content_subtype() == "html":
        body = _html_to_text(body)
    if not body.strip():
        raise ExtractionError("Email has no text body")
    return "\n".join(header) + "\n\n" + body.rstrip("\n") + "\n"


def _extract_epub(path: str) -> str:
    _require("ebooklib", "ebooklib")
    import ebooklib
    from ebooklib import epub

    # Same contract concern as _extract_msg: ebooklib raises its own errors on
    # a malformed archive, which must not escape as anything but ExtractionError.
    try:
        book = epub.read_epub(path)
        items = list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))
    except ExtractionError:
        raise
    except Exception as exc:  # noqa: BLE001 - third-party exception tree
        raise ExtractionError(f"Not a readable EPUB: {exc}") from exc

    chunks = [_html_to_text(item.get_content().decode("utf-8", errors="replace")) for item in items]
    text = "\n\n".join(chunk for chunk in chunks if chunk.strip())
    if not text.strip():
        raise ExtractionError("EPUB contains no extractable text")
    return text.rstrip("\n") + "\n"


def _html_to_text(markup: str) -> str:
    _require("bs4", "beautifulsoup4")
    from bs4 import BeautifulSoup

    # lxml is the declared parser; fall back to the stdlib one so a partial
    # install still renders HTML rather than failing outright.
    try:
        soup = BeautifulSoup(markup, "lxml")
    except Exception:  # noqa: BLE001 - lxml missing/broken
        soup = BeautifulSoup(markup, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return re.sub(r"\n{3,}", "\n\n", soup.get_text("\n")).strip()


def _extract_html(path: str) -> str:
    try:
        markup = Path(path).read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        raise ExtractionError(str(exc)) from exc
    text = _html_to_text(markup)
    if not text.strip():
        raise ExtractionError("HTML contains no extractable text")
    return text.rstrip("\n") + "\n"


_PDF_OCR_DPI = 200


def _ocr_tesseract(image) -> str:
    """OCR engine #1. Needs the `tesseract` binary the installer provisions."""
    try:
        import pytesseract
    except ImportError as exc:
        raise ExtractionError(f"pytesseract is not installed: {exc}") from exc
    try:
        return pytesseract.image_to_string(image)
    except pytesseract.pytesseract.TesseractNotFoundError as exc:
        raise ExtractionError(
            "the tesseract binary is not on PATH (re-run the Pocura installer, "
            "or install tesseract manually)"
        ) from exc


_rapidocr_engine = None


def _ocr_rapidocr(image) -> str:
    """OCR engine #2 (ONNX). Pure-wheel, so it works where the tesseract
    binary is missing entirely -- which is the main reason it's here.

    The engine loads its models on construction (slow, ~100MB), so it's built
    once and cached for the process.
    """
    global _rapidocr_engine
    try:
        from rapidocr_onnxruntime import RapidOCR
    except ImportError as exc:
        raise ExtractionError(f"rapidocr-onnxruntime is not installed: {exc}") from exc

    import numpy as np

    if _rapidocr_engine is None:
        _rapidocr_engine = RapidOCR()
    # RapidOCR wants an array; hand it RGB so mode-specific quirks (P, LA,
    # RGBA from a PDF pixmap) can't reach it.
    result, _elapsed = _rapidocr_engine(np.array(image.convert("RGB")))
    if not result:
        return ""
    return "\n".join(line[1] for line in result)


def _ocr(image) -> str:
    """Run the OCR engines in order, first non-empty result wins.

    Returns "" (rather than raising) when every engine is unavailable or finds
    nothing, so a PDF page with no text simply reads as empty instead of
    failing the whole document -- callers decide whether empty is fatal.
    """
    for name, engine in (("tesseract", _ocr_tesseract), ("rapidocr", _ocr_rapidocr)):
        try:
            text = engine(image)
        except ExtractionError as exc:
            logger.debug("OCR engine %s unavailable: %s", name, exc)
            continue
        except Exception:  # noqa: BLE001 - a broken engine must not mask the other
            logger.debug("OCR engine %s failed", name, exc_info=True)
            continue
        if text and text.strip():
            return text
    return ""


_PAGE_EMPTY = "(no extractable text)"


def _pages_to_text(pages: list[str], label: str) -> str:
    out: list[str] = []
    for idx, text in enumerate(pages, start=1):
        out.append(f"# ── Page {idx} ──")
        out.append(text.strip() or _PAGE_EMPTY)
        out.append("")
    if not any(text.strip() for text in pages):
        raise ExtractionError(f"{label} produced no text")
    return "\n".join(out).rstrip("\n") + "\n"


def _extract_pdf_pymupdf(path: str) -> str:
    """Primary PDF path: text layer per page, OCR for pages that have none.

    This is the only extractor that can rasterize, so it's the only one that
    handles scanned PDFs -- the others are text-layer-only and will come back
    empty on a scan, which is exactly what the chain is for.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise ExtractionError(f"PyMuPDF is not installed: {exc}") from exc
    from PIL import Image

    try:
        doc = fitz.open(path)
    except Exception as exc:
        raise ExtractionError(f"PyMuPDF could not open the file: {exc}") from exc

    pages: list[str] = []
    try:
        for page in doc:
            text = page.get_text().strip()
            if not text:
                pixmap = page.get_pixmap(dpi=_PDF_OCR_DPI)
                mode = "RGBA" if pixmap.alpha else "RGB"
                image = Image.frombytes(mode, (pixmap.width, pixmap.height), pixmap.samples)
                text = _ocr(image).strip()
            pages.append(text)
    finally:
        doc.close()
    return _pages_to_text(pages, "PyMuPDF")


def _extract_pdf_pdfplumber(path: str) -> str:
    """Layout/table-aware text layer. No rasterization, so scans read empty."""
    _require("pdfplumber", "pdfplumber")
    import pdfplumber

    with pdfplumber.open(path) as pdf:
        pages = [(page.extract_text() or "").strip() for page in pdf.pages]
    return _pages_to_text(pages, "pdfplumber")


def _extract_pdf_pypdf(path: str) -> str:
    """Text layer via pypdf — handles some odd xref tables PyMuPDF rejects."""
    _require("pypdf", "pypdf")
    import pypdf

    reader = pypdf.PdfReader(path)
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    return _pages_to_text(pages, "pypdf")


_java_available: bool | None = None


def _has_java() -> bool:
    global _java_available
    if _java_available is None:
        _java_available = shutil.which("java") is not None
    return _java_available


def _extract_pdf_opendataloader(path: str) -> str:
    """Structured extraction via opendataloader-pdf.

    It's a thin Python wrapper over a Java CLI, so it is useless without a
    JRE 11+ — which we deliberately do not install. Probe for `java` first and
    bail with a plain message rather than letting a Java stack trace surface.
    """
    if not _has_java():
        raise ExtractionError("needs a Java 11+ runtime, which is not installed")
    opendataloader_pdf = _require("opendataloader_pdf", "opendataloader-pdf")

    text = opendataloader_pdf.run([path], output_format=["markdown"])
    if isinstance(text, (list, tuple)):
        text = "\n\n".join(str(chunk) for chunk in text)
    text = str(text or "")
    if not text.strip():
        raise ExtractionError("opendataloader produced no text")
    return text.rstrip("\n") + "\n"


def _extract_pdf(path: str) -> str:
    return _first_non_empty(
        path,
        [
            ("pymupdf+ocr", _extract_pdf_pymupdf),
            ("pdfplumber", _extract_pdf_pdfplumber),
            ("pypdf", _extract_pdf_pypdf),
            ("opendataloader", _extract_pdf_opendataloader),
        ],
        "PDF",
    )


def _extract_image(path: str) -> str:
    from PIL import Image

    try:
        with Image.open(path) as image:
            image.load()
            text = _ocr(image)
    except Exception as exc:  # noqa: BLE001 - Pillow raises a wide range here
        raise ExtractionError(f"Not a readable image: {exc}") from exc

    if not text.strip():
        raise ExtractionError(
            "No OCR-detectable text in this image (no OCR engine available, or "
            "the image contains no text)"
        )
    return text.rstrip("\n") + "\n"


def _shared_strings(zf: zipfile.ZipFile, names: set[str]) -> list[str]:
    if "xl/sharedStrings.xml" not in names:
        return []
    try:
        root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
    except ET.ParseError:
        return []
    s = f"{{{_NS_S}}}"
    return ["".join(t.text or "" for t in item.iter(f"{s}t")) for item in root.iter(f"{s}si")]


def _workbook_sheets(zf: zipfile.ZipFile) -> list[tuple[str, str, str]]:
    root = _zip_xml(zf, "xl/workbook.xml")
    s, r = f"{{{_NS_S}}}", f"{{{_NS_REL}}}"
    return [
        (sheet.get("name", "Sheet"), sheet.get("state", "visible"), sheet.get(f"{r}id", ""))
        for sheet in root.iter(f"{s}sheet")
    ]


def _workbook_rels(zf: zipfile.ZipFile, names: set[str]) -> dict[str, str]:
    rels_path = "xl/_rels/workbook.xml.rels"
    if rels_path not in names:
        return {}
    try:
        root = ET.fromstring(zf.read(rels_path))
    except ET.ParseError:
        return {}
    rel_tag = f"{{{_NS_PKG_REL}}}Relationship"
    return {rel.get("Id", ""): rel.get("Target", "") for rel in root.iter(rel_tag) if rel.get("Id")}


def _sheet_part(target: str) -> str:
    target = target.lstrip("/")
    return posixpath.normpath(target if target.startswith("xl/") else f"xl/{target}")


def _col_index(ref: str) -> int:
    idx = 0
    for ch in ref:
        if not ch.isalpha():
            break
        idx = idx * 26 + ord(ch.upper()) - ord("A") + 1
    return max(idx - 1, 0)


def _sheet_rows(xml_bytes: bytes, shared: list[str]) -> list[list[str]]:
    root = ET.fromstring(xml_bytes)
    s = f"{{{_NS_S}}}"
    rows: list[list[str]] = []
    for row in root.iter(f"{s}row"):
        if len(rows) >= _MAX_XLSX_ROWS_PER_SHEET:
            break
        cells: dict[int, str] = {}
        max_col = -1
        for cell in row.iter(f"{s}c"):
            col = _col_index(cell.get("r", "")) if cell.get("r") else max_col + 1
            if col >= _MAX_XLSX_COLS:
                continue
            cells[col] = _cell_value(cell, shared, s)
            max_col = max(max_col, col)
        rows.append([cells.get(i, "") for i in range(max_col + 1)] if max_col >= 0 else [])
    while rows and not any(value.strip() for value in rows[-1]):
        rows.pop()
    return rows


def _cell_value(cell: ET.Element, shared: list[str], s: str) -> str:
    value = cell.findtext(f"{s}v") or ""
    typ = cell.get("t", "")
    if typ == "s":
        try:
            return shared[int(value)]
        except (ValueError, IndexError):
            return ""
    if typ == "inlineStr":
        inline = cell.find(f"{s}is")
        return "" if inline is None else "".join(t.text or "" for t in inline.iter(f"{s}t"))
    if typ == "b":
        return "TRUE" if value.strip() in {"1", "true", "TRUE"} else "FALSE"
    if typ == "e":
        return value or "#ERROR"
    return value
